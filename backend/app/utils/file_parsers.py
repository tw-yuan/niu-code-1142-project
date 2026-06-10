import base64
import io
from pathlib import Path


# ── PDF ──────────────────────────────────────────────────────────────────────

def _pdf_text_per_page(doc) -> list[int]:
    """回傳每頁可萃取的字元數（含空白）"""
    return [len(doc[i].get_text().strip()) for i in range(len(doc))]


def _page_to_jpeg_b64(page, zoom: float, quality: int) -> str:
    import fitz
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
    jpeg_bytes = pix.tobytes("jpeg", jpg_quality=quality)
    return base64.b64encode(jpeg_bytes).decode()


def _page_has_visual_blocks(page) -> bool:
    """判斷頁面是否含圖片、向量圖形或繪圖元素。"""
    try:
        if page.get_images(full=True):
            return True
    except Exception:
        pass
    try:
        blocks = page.get_text("dict").get("blocks", [])
        if any(block.get("type") == 1 for block in blocks):
            return True
    except Exception:
        pass
    try:
        return bool(page.get_drawings())
    except Exception:
        return False


def parse_pdf_text(file_bytes: bytes) -> str:
    """純文字 PDF：用 pymupdf4llm 輸出 Markdown。"""
    import fitz
    import pymupdf4llm
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    return pymupdf4llm.to_markdown(doc)


async def parse_pdf_vision(file_bytes: bytes, force_all_pages: bool = False) -> str:
    """
    掃描版 PDF（或混合型）：
    - 低文字頁、圖片頁、圖形頁轉 JPEG 送視覺模型
    - 其他純文字頁用 pymupdf4llm
    預設每次只送 1 頁，避免 OpenRouter 或其他 OpenAI-compatible provider 回 413。
    """
    import fitz
    import pymupdf4llm
    from openai import AsyncOpenAI
    from app.config import settings

    if settings.demo_mode or not settings.openai_compatible_api_key:
        text = parse_pdf_text(file_bytes).strip()
        warning = (
            "注意：這份 PDF 含掃描頁、圖片或圖形，但目前未啟用 vision 解析。"
            "請設定 OPENAI_COMPATIBLE_API_KEY、使用支援圖片的 VISION_MODEL，並確認 DEMO_MODE=false，才能完整讀取圖片型 PDF、圖表與版面。"
        )
        return f"{text}\n\n---\n\n{warning}" if text else warning

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chars_per_page = _pdf_text_per_page(doc)
    threshold = settings.pdf_text_threshold

    client = AsyncOpenAI(
        base_url=settings.openai_compatible_base_url,
        api_key=settings.openai_compatible_api_key or "none",
    )

    batch_size = max(1, settings.pdf_vision_batch_size)
    zoom = max(0.8, settings.pdf_image_zoom)
    quality = min(95, max(40, settings.pdf_image_jpeg_quality))
    results: list[str] = []

    for batch_start in range(0, len(doc), batch_size):
        batch_end = min(batch_start + batch_size, len(doc))
        image_pages: list[int] = []
        text_pages: list[int] = []

        for i in range(batch_start, batch_end):
            needs_vision = (
                force_all_pages
                or chars_per_page[i] < threshold
                or _page_has_visual_blocks(doc[i])
            )
            if needs_vision:
                image_pages.append(i)
            else:
                text_pages.append(i)

        # 文字頁 → pymupdf4llm（只處理這批的頁次）
        if text_pages:
            # pymupdf4llm 支援指定頁碼列表
            md = pymupdf4llm.to_markdown(doc, pages=text_pages)
            results.append(md)

        # 圖片頁 → 視覺模型
        if image_pages:
            content: list[dict] = []
            for i in image_pages:
                b64 = _page_to_jpeg_b64(doc[i], zoom=zoom, quality=quality)
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{b64}",
                        "detail": "high",
                    },
                })
            page_range = f"第 {image_pages[0]+1}～{image_pages[-1]+1} 頁"
            content.append({
                "type": "text",
                "text": (
                    f"以上是 PDF {page_range} 的圖片。"
                    "請將頁面中可見文字、圖表、流程圖、表格、圖片註解與版面關係完整轉換為繁體中文 Markdown。"
                    "若有圖表或流程圖，請用文字描述其重點與關係。保留標題層級、粗體、表格。"
                    "直接輸出 Markdown，不需額外說明。"
                ),
            })
            resp = await client.chat.completions.create(
                model=settings.vision_model,
                messages=[{"role": "user", "content": content}],
                max_tokens=4096,
            )
            results.append(resp.choices[0].message.content or "")

    return "\n\n---\n\n".join(filter(None, results))


def _is_scanned_pdf(file_bytes: bytes, threshold: int) -> bool:
    """超過一半頁面文字不足閾值，判定為掃描版。"""
    import fitz
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    if len(doc) == 0:
        return False
    sparse = sum(1 for c in _pdf_text_per_page(doc) if c < threshold)
    return sparse / len(doc) > 0.5


def _pdf_needs_visual_parse(file_bytes: bytes, threshold: int) -> bool:
    """只要 PDF 有掃描頁、圖片頁或圖形頁，就啟用混合視覺解析。"""
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    if len(doc) == 0:
        return False
    chars_per_page = _pdf_text_per_page(doc)
    if sum(1 for c in chars_per_page if c < threshold) / len(doc) > 0.5:
        return True
    return any(_page_has_visual_blocks(doc[i]) for i in range(len(doc)))


# ── 其他格式 ─────────────────────────────────────────────────────────────────

def parse_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def parse_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[str] = []

    def collect_shape(shape) -> list[str]:
        parts: list[str] = []
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                parts.extend(collect_shape(child))
            return parts
        try:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
        except Exception:
            pass
        try:
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                if rows:
                    parts.append("\n".join(rows))
        except Exception:
            pass
        try:
            if getattr(shape, "has_chart", False):
                title = ""
                if shape.chart.has_title and shape.chart.chart_title.text_frame:
                    title = shape.chart.chart_title.text_frame.text.strip()
                parts.append(f"[圖表{f'：{title}' if title else ''}]")
        except Exception:
            pass
        try:
            if getattr(shape, "shape_type", None) and "PICTURE" in str(shape.shape_type):
                parts.append("[圖片]")
        except Exception:
            pass
        return parts

    for slide_index, slide in enumerate(prs.slides, start=1):
        parts = [f"## 投影片 {slide_index}"]
        for shape in slide.shapes:
            parts.extend(collect_shape(shape))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"### 備註\n{notes}")
        except Exception:
            pass
        if len(parts) == 1:
            parts.append("[此投影片沒有可直接抽取的文字，可能主要由圖片、圖形或嵌入物件構成。]")
        slides.append("\n\n".join(parts))
    return "\n\n---\n\n".join(slides)


def parse_text(file_bytes: bytes) -> str:
    for enc in ("utf-8", "big5", "gbk", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except (UnicodeDecodeError, ValueError):
            continue
    return file_bytes.decode("utf-8", errors="replace")


async def parse_image_vision(filename: str, file_bytes: bytes) -> str:
    from openai import AsyncOpenAI
    from app.config import settings

    if settings.demo_mode or not settings.openai_compatible_api_key:
        return f"這是一張課程圖片或講義截圖（{filename}）。示範模式未呼叫視覺模型，請設定 VISION_MODEL API key 取得完整辨識內容。"

    ext = Path(filename).suffix.lower().lstrip(".") or "png"
    mime = "jpeg" if ext in ("jpg", "jpeg") else ext
    b64 = base64.b64encode(file_bytes).decode()
    client = AsyncOpenAI(
        base_url=settings.openai_compatible_base_url,
        api_key=settings.openai_compatible_api_key or "none",
    )
    resp = await client.chat.completions.create(
        model=settings.vision_model,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/{mime};base64,{b64}",
                            "detail": "high",
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "請將這張課程講義、白板、投影片或筆記圖片完整轉成繁體中文 Markdown。"
                            "保留標題、條列、表格與重要符號。直接輸出 Markdown，不需額外說明。"
                        ),
                    },
                ],
            }
        ],
        max_tokens=4096,
    )
    return resp.choices[0].message.content or ""


# ── 統一入口 ─────────────────────────────────────────────────────────────────

async def parse_file_async(filename: str, file_bytes: bytes) -> str:
    """非同步版本，PDF 視情況走視覺路徑。"""
    from app.config import settings
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        strategy = settings.pdf_vision_strategy.lower()
        if strategy == "always":
            return await parse_pdf_vision(file_bytes, force_all_pages=True)
        if strategy != "never" and _pdf_needs_visual_parse(file_bytes, settings.pdf_text_threshold):
            return await parse_pdf_vision(file_bytes)
        return parse_pdf_text(file_bytes)
    if ext == ".docx":
        return parse_docx(file_bytes)
    if ext == ".pptx":
        return parse_pptx(file_bytes)
    if ext in {".jpg", ".jpeg", ".png", ".webp"}:
        return await parse_image_vision(filename, file_bytes)
    return parse_text(file_bytes)


def count_tokens(text: str) -> int:
    import tiktoken
    enc = tiktoken.get_encoding("cl100k_base")
    return len(enc.encode(text))
