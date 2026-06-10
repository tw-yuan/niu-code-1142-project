from io import BytesIO

from docx import Document
from pptx import Presentation

from app.services.chat_service import _demo_response, _quiz_metadata
from app.services.direction_service import fallback_dynamic_directions
from app.utils.file_parsers import _docx_is_plain_text, parse_pptx
from app.services.rag_service import _source_label


def test_fallback_dynamic_directions_returns_two_to_three_items():
    directions = fallback_dynamic_directions("資料結構 stack queue tree algorithm", "資料結構")

    assert 2 <= len(directions) <= 3
    assert all(item["is_dynamic"] for item in directions)
    assert {item["key"] for item in directions}


def test_demo_quiz_response_can_be_stored_as_quiz_metadata():
    response = _demo_response("quiz", "自我測驗", "請出題", "## Stack\n先進後出資料結構")
    metadata = _quiz_metadata("quiz", "請出題", response)

    assert metadata is not None
    assert metadata["kind"] == "quiz"
    assert metadata["status"] == "generated"
    assert metadata["question_count"] == 3


def test_quiz_metadata_extracts_score_when_available():
    metadata = _quiz_metadata("quiz", "我的答案", "批改結果：80 分。觀念大致正確。")

    assert metadata is not None
    assert metadata["status"] == "graded"
    assert metadata["score"] == 80


def test_source_label_uses_first_meaningful_line():
    assert _source_label("\n# 第三章 樹狀結構\n內容", 0) == "第三章 樹狀結構"


def test_source_label_falls_back_to_chunk_number():
    assert _source_label("\n\n", 2) == "片段 3"


def test_parse_pptx_extracts_text_table_and_empty_slide_marker():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "資料結構"
    table_shape = slide.shapes.add_table(2, 2, 0, 1000000, 4000000, 1000000)
    table = table_shape.table
    table.cell(0, 0).text = "結構"
    table.cell(0, 1).text = "特性"
    table.cell(1, 0).text = "Stack"
    table.cell(1, 1).text = "LIFO"
    prs.slides.add_slide(prs.slide_layouts[6])

    buf = BytesIO()
    prs.save(buf)

    parsed = parse_pptx(buf.getvalue())

    assert "資料結構" in parsed
    assert "Stack | LIFO" in parsed
    assert "沒有可直接抽取的文字" in parsed


def test_docx_plain_text_detection_rejects_tables():
    plain = Document()
    plain.add_paragraph("這是一份純文字講義")
    plain_buf = BytesIO()
    plain.save(plain_buf)

    with_table = Document()
    table = with_table.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "欄位"
    table.cell(0, 1).text = "內容"
    table_buf = BytesIO()
    with_table.save(table_buf)

    assert _docx_is_plain_text(plain_buf.getvalue()) is True
    assert _docx_is_plain_text(table_buf.getvalue()) is False
